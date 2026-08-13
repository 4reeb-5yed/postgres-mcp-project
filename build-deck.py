"""Builds the Claude DB Connect showcase deck as an editable .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

OUT = "Claude-DB-Connect-Showcase.pptx"
SHOTS = "docs/screenshots"

# ---------- design system ----------
INK       = RGBColor(0x16, 0x18, 0x1D)
INK_SOFT  = RGBColor(0x24, 0x27, 0x2E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0x1A, 0x1D, 0x21)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
FAINT     = RGBColor(0x9C, 0xA3, 0xAF)
RULE      = RGBColor(0xE5, 0xE7, 0xEB)
PANEL     = RGBColor(0xF7, 0xF8, 0xFA)
CODE_BG   = RGBColor(0x1E, 0x21, 0x28)
CODE_FG   = RGBColor(0xE6, 0xE9, 0xEF)
CODE_KEY  = RGBColor(0x7EE, 0x0, 0x0) if False else RGBColor(0x8B, 0xD4, 0x9C)
CODE_DIM  = RGBColor(0x9A, 0xA3, 0xB2)

BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x15, 0x80, 0x3D)
AMBER  = RGBColor(0xB4, 0x53, 0x09)
RED    = RGBColor(0xB9, 0x1C, 0x1C)

PG = RGBColor(0x33, 0x67, 0x91)
MY = RGBColor(0xD9, 0x77, 0x06)
MS = RGBColor(0xA4, 0x26, 0x2C)
SQ = RGBColor(0x0F, 0x6F, 0xB5)

H_FONT = "Segoe UI"
C_FONT = "Consolas"

SW, SH = 13.333, 7.5
ML, MR = 0.72, 0.72
CW = SW - ML - MR

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK if dark else WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    _page["n"] += 1
    return s


def txt(s, x, y, w, h, runs, size=14, bold=False, color=TEXT, font=H_FONT,
        align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP, space_after=6):
    """runs: str, or list of paragraphs; each paragraph str or list of (text, opts)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        chunks = para if isinstance(para, list) else [(para, {})]
        for t, o in chunks:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = o.get("font", font)
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", bold)
            f.italic = o.get("italic", False)
            f.color.rgb = o.get("color", color)
    return tb


def rect(s, x, y, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def chip(s, y, label, color):
    rect(s, ML, y, 0.17, 0.17, fill=color)
    txt(s, ML + 0.3, y - 0.035, 8.0, 0.3, label.upper(), size=10.5, bold=True,
        color=color, spacing=1.0)


def head(s, y_chip, label, color, title, sub=None, tsize=30):
    chip(s, y_chip, label, color)
    txt(s, ML, y_chip + 0.32, CW, 0.62, title, size=tsize, bold=True, color=TEXT, spacing=0.95)
    yy = y_chip + 0.32 + (0.66 if tsize >= 28 else 0.56)
    if sub:
        txt(s, ML, yy, CW - 0.4, 0.4, sub, size=13.5, color=MUTED, spacing=1.15)
        yy += 0.44
    rect(s, ML, yy + 0.07, 1.5, 0.035, fill=color)
    rect(s, ML + 1.5, yy + 0.07, CW - 1.5, 0.035, fill=RULE)
    return yy + 0.34


def callout(s, y, label, body, color, h=0.72, w=None):
    w = w or CW
    rect(s, ML, y, w, h, fill=PANEL)
    rect(s, ML, y, 0.055, h, fill=color)
    txt(s, ML + 0.26, y + 0.13, w - 0.5, h - 0.2,
        [[(label + "  ", {"bold": True, "color": color, "size": 12.5}),
          (body, {"color": TEXT, "size": 12.5})]], spacing=1.2)


def code(s, x, y, w, h, lines, size=11.5, title=None):
    if title:
        txt(s, x, y - 0.28, w, 0.25, title, size=10.5, bold=True, color=MUTED)
    rect(s, x, y, w, h, fill=CODE_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.035)
    paras = []
    for ln in lines:
        if isinstance(ln, tuple):
            paras.append([(ln[0], {"font": C_FONT, "size": size, "color": ln[1]})])
        else:
            paras.append([(ln, {"font": C_FONT, "size": size, "color": CODE_FG})])
    txt(s, x + 0.24, y + 0.18, w - 0.48, h - 0.36, paras, spacing=1.28, space_after=0)


def fit_image(s, path, bx, by, bw, bh, border=True):
    with Image.open(path) as im:
        pw, ph = im.size
    ar = ph / pw
    w = bw
    h = bw * ar
    if h > bh:
        h = bh
        w = bh / ar
    x = bx + (bw - w) / 2
    y = by + (bh - h) / 2
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if border:
        try:
            pic.line.color.rgb = RULE
            pic.line.width = Pt(0.75)
        except Exception:
            pass
    return pic


def footer(s, label, dark=False):
    n = _page["n"]
    txt(s, ML, SH - 0.46, 8.0, 0.26, label, size=9.5,
        color=(RGBColor(0x6E, 0x74, 0x80) if dark else FAINT), spacing=1.0)
    txt(s, SW - MR - 0.8, SH - 0.46, 0.8, 0.26, str(n), size=9.5,
        color=(RGBColor(0x6E, 0x74, 0x80) if dark else FAINT), align=PP_ALIGN.RIGHT)





# ============================== 1 · TITLE ==============================
s = slide(dark=True)
rect(s, 0, 0, 0.16, SH, fill=BLUE)
for i, c in enumerate([PG, MY, MS, SQ]):
    rect(s, ML + i * 0.42, 1.55, 0.22, 0.22, fill=c, shape=MSO_SHAPE.OVAL)
txt(s, ML, 2.1, 11.4, 1.7,
    [[("Claude DB Connect", {})],
     [("One assistant, four local databases", {"color": BLUE})]],
    size=42, bold=True, color=WHITE, spacing=0.92)
txt(s, ML, 4.0, 9.8, 1.0,
    "Claude Desktop wired to PostgreSQL, MySQL, MS SQL Server and SQLite over the Model "
    "Context Protocol — asking questions in plain English and getting answers from real SQL "
    "run against live local databases.",
    size=15, color=RGBColor(0xB6, 0xBD, 0xC8), spacing=1.35)
rect(s, ML, 5.35, 2.2, 0.045, fill=BLUE)
txt(s, ML, 5.62, 11.0, 0.4,
    "What it is  ·  how it works  ·  a full MS SQL Server walkthrough  ·  what broke along the way",
    size=11.5, color=RGBColor(0x7C, 0x84, 0x92), spacing=1.2)
footer(s, "Claude DB Connect", dark=True)

# ============================== 2 · WHAT THIS IS ==============================
s = slide()
y = head(s, 0.46, "The project", BLUE, "What this is, and what you can do with it",
         "A working integration, plus the debugging trail that made it work")
txt(s, ML, y, CW, 0.85,
    [[("Claude Desktop can normally only talk ", {}),
      ("about", {"italic": True}),
      (" your data. This project gives it a live connection to four different local databases, "
       "so it reads the actual schema, writes real SQL, runs it, and answers from the rows that "
       "come back — with the tool call visible every time.", {})]],
    size=14.5, color=TEXT, spacing=1.4)

cards = [
    ("Ask in plain English", BLUE,
     "\u201cWhat columns does the employees table have?\u201d becomes a real schema query. No SQL "
     "written by hand, no connection string pasted into a chat."),
    ("Four databases, side by side", GREEN,
     "PostgreSQL, MySQL, MS SQL Server and SQLite all connected at once — so a prompt can name "
     "which engine it should hit."),
    ("It generates working code", AMBER,
     "Ask for a Python script that charts the data and it pulls the live rows first, then writes "
     "code against the real schema."),
]
cy = y + 1.05
cwid = (CW - 0.36) / 3
for i, (t, c, b) in enumerate(cards):
    cx = ML + i * (cwid + 0.18)
    rect(s, cx, cy, cwid, 1.95, fill=PANEL)
    rect(s, cx, cy, cwid, 0.05, fill=c)
    txt(s, cx + 0.24, cy + 0.28, cwid - 0.48, 0.3, t, size=13.5, bold=True, color=TEXT)
    txt(s, cx + 0.24, cy + 0.68, cwid - 0.48, 1.1, b, size=11.5, color=MUTED, spacing=1.3)

callout(s, cy + 2.2,
        "Worth saying plainly:",
        "not one of the four worked by following the official guide as written. The integration is "
        "the deliverable; the debugging trail is the interesting part, and it is documented per "
        "database in the repo.", BLUE, h=0.8)
footer(s, "The project")

# ============================== 3 · ARCHITECTURE ==============================
s = slide()
y = head(s, 0.46, "How it works", GREEN, "The architecture, end to end",
         "Every one of the four integrations is the same three-hop path")

bw_, bh_, gap_ = 3.35, 1.5, 0.72
by_ = y + 0.5
boxes = [
    ("Claude Desktop", "the MCP client", BLUE),
    ("MCP server", "translates tool calls to SQL", GREEN),
    ("Your database", "running locally", AMBER),
]
for i, (t, sub, c) in enumerate(boxes):
    bx_ = ML + i * (bw_ + gap_)
    rect(s, bx_, by_, bw_, bh_, fill=WHITE, line=RULE, lw=1.0)
    rect(s, bx_, by_, 0.05, bh_, fill=c)
    txt(s, bx_ + 0.26, by_ + 0.34, bw_ - 0.5, 0.3, t, size=15, bold=True, color=TEXT)
    txt(s, bx_ + 0.26, by_ + 0.74, bw_ - 0.5, 0.5, sub, size=11.5, color=MUTED, spacing=1.2)
    if i < 2:
        ax = bx_ + bw_ + 0.08
        ar_ = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(by_ + bh_ / 2 - 0.13),
                                 Inches(gap_ - 0.16), Inches(0.26))
        ar_.fill.solid(); ar_.fill.fore_color.rgb = RGBColor(0xC7, 0xCD, 0xD6)
        ar_.line.fill.background(); ar_.shadow.inherit = False
txt(s, ML + bw_ + 0.06, by_ + bh_ + 0.12, gap_ + 0.6, 0.3, "stdio", size=10.5, bold=True,
    color=GREEN, align=PP_ALIGN.CENTER)
txt(s, ML + 2 * (bw_ + gap_) - gap_ - 0.06, by_ + bh_ + 0.12, gap_ + 0.6, 0.3, "SQL", size=10.5,
    bold=True, color=AMBER, align=PP_ALIGN.CENTER)

txt(s, ML, by_ + bh_ + 0.62, CW, 1.5,
    [[("The middle hop is the whole story. ", {"bold": True}),
      ("Claude Desktop launches each MCP server as a local child process and talks to it over "
       "stdio \u2014 plain stdin/stdout, no network port, no HTTP. The server exposes database "
       "operations as callable tools, so Claude can list tables, read a schema, and run a query "
       "as discrete tool calls.", {})],
     [("This single constraint is what invalidated the setup the original guide recommended: the "
       "image it named runs an HTTP server, which Claude Desktop's local MCP integration cannot "
       "speak to at all.", {"color": MUTED, "size": 13})]],
    size=14, color=TEXT, spacing=1.4, space_after=10)

callout(s, 6.35, "Nothing is hosted:",
        "the databases, the MCP servers and the client all run on one machine. No cloud database, "
        "no exposed port, no shared credentials.", GREEN, h=0.6)
footer(s, "How it works")


# ============================== 5 · THE FOUR AT A GLANCE ==============================
s = slide()
y = head(s, 0.46, "How it works", GREEN, "The four integrations at a glance",
         "Same architecture every time — what changes is how the server is packaged and where the data lives")

rows = [
    ("Database", "MCP server", "Launched via", "Where the data lives", None),
    ("PostgreSQL", "postgres-mcp", "uvx, with mcp<2.0.0 pinned", "Docker container, port 5433", PG),
    ("MySQL", "mysql-mcp-server", "uvx --from", "Docker container, port 3306", MY),
    ("MS SQL Server", "sql-mcp-server", "pip install + ODBC Driver 17", "Docker container, port 1433", MS),
    ("SQLite", "mcp-server-sqlite", "uvx", "A local .db file — no server", SQ),
]
tw = [2.55, 2.95, 3.55, 3.84]
th = 0.56
gt = s.shapes.add_table(len(rows), 4, Inches(ML), Inches(y + 0.12),
                        Inches(CW), Inches(th * len(rows))).table
gt.first_row = False
for i, w in enumerate(tw):
    gt.columns[i].width = Inches(w)
for ri, r in enumerate(rows):
    gt.rows[ri].height = Inches(th)
    for ci in range(4):
        cell = gt.cell(ri, ci)
        cell.text = ""
        cell.margin_left = Inches(0.18)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.07)
        cell.margin_bottom = Inches(0.07)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = INK
        else:
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else PANEL
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = r[ci]
        f = run.font
        f.name = H_FONT
        f.size = Pt(11.5 if ri else 11)
        f.bold = (ri == 0) or (ci == 0)
        if ri == 0:
            f.color.rgb = WHITE
        elif ci == 0:
            f.color.rgb = r[4]
        elif ci == 1:
            f.name = C_FONT
            f.size = Pt(11)
            f.color.rgb = TEXT
        else:
            f.color.rgb = MUTED

callout(s, y + 0.12 + th * len(rows) + 0.28, "Full step-by-step for each:",
        "docs/setup/postgres.md \u00b7 mysql.md \u00b7 mssql.md \u00b7 sqlite.md \u2014 every command, config block "
        "and gotcha, per database. This deck walks through MS SQL Server in full and points to the "
        "rest.", BLUE, h=0.78)
footer(s, "How it works")

# ============================== 6 · WHY FOUR ==============================
s = slide()
y = head(s, 0.46, "How it works", GREEN, "Why four, and what actually differs",
         "One integration shows the pattern; four expose where the real friction is")

items = [
    ("Same transport, different packaging", GREEN,
     "Three of the four launch through uvx, which fetches and runs the Python server on demand. "
     "MS SQL Server needs a real pip install instead, because its server expects a system-level "
     "database driver to already be present."),
    ("One needs a host-level dependency", MS,
     "ODBC Driver 17 for SQL Server has to be installed on the machine itself, with admin "
     "elevation. Without it pyodbc cannot open a connection, and the MCP server fails at startup "
     "rather than at query time."),
    ("One has no server at all", SQ,
     "SQLite is just a file on disk. No container to start, no port to bind, no password \u2014 the "
     "config simply points at a .db path, which makes it the shortest setup of the four by far."),
    ("The same upstream break hit two of them", RED,
     "postgres-mcp and sql-mcp-server both crashed on a breaking mcp 2.0.0 release. Same root "
     "cause, two different fixes: pin at invocation for uvx, pin at install for pip."),
]
iy = y + 0.28
for t, c, b in items:
    rect(s, ML, iy, 0.05, 0.92, fill=c)
    txt(s, ML + 0.26, iy + 0.02, CW - 0.4, 0.28, t, size=13.5, bold=True, color=TEXT)
    txt(s, ML + 0.26, iy + 0.34, CW - 0.4, 0.6, b, size=11.5, color=MUTED, spacing=1.28)
    iy += 1.12

callout(s, iy + 0.02, "The takeaway:",
        "\u201cconnect an AI assistant to a database\u201d is one sentence and four genuinely different "
        "problems.", GREEN, h=0.52)
footer(s, "How it works")

# ============================== 7 · DIVIDER ==============================
s = slide(dark=True)
rect(s, 0, 0, 0.16, SH, fill=MS)
txt(s, ML, 1.95, 4.0, 0.35, "THE WALKTHROUGH", size=12, bold=True, color=MS)
txt(s, ML, 2.38, 10.6, 1.0, "Setting up MS SQL Server", size=36, bold=True, color=WHITE,
    spacing=0.95)
rect(s, ML, 3.5, 2.0, 0.045, fill=MS)
txt(s, ML, 3.82, 9.6, 1.5,
    "Taking one of the four end to end \u2014 a host-level ODBC driver, a pip install instead of uvx, "
    "and a dependency pin to get past an upstream breaking change. Four steps, then what it can do.",
    size=15, color=RGBColor(0xB6, 0xBD, 0xC8), spacing=1.35)
rect(s, ML, 5.3, 9.6, 0.82, fill=INK_SOFT)
rect(s, ML, 5.3, 0.05, 0.82, fill=MS)
txt(s, ML + 0.3, 5.47, 9.1, 0.55,
    [[("The other three  ", {"bold": True, "color": MS, "size": 12}),
      ("follow the same four-step shape \u2014 their full guides live in docs/setup/ in the repo.",
       {"color": RGBColor(0xB6, 0xBD, 0xC8), "size": 12})]], spacing=1.25)
footer(s, "MS SQL Server", dark=True)

# ============================== 8 · STEP 1 ==============================
s = slide()
y = head(s, 0.44, "Step 1 of 4 \u00b7 MS SQL Server", MS, "Start the server, install the driver",
         "Two things have to be in place before any MCP server can connect", tsize=28)
txt(s, ML, y, CW, 0.72,
    [[("The container is standard. ", {"bold": True}),
      ("The part that is easy to miss is the second command: ODBC Driver 17 is a host-level "
       "install, not a Python package, and it needs admin elevation. Without it the MCP server "
       "starts and then fails to connect, with an error that points at the driver rather than "
       "your password.", {})]],
    size=13.5, spacing=1.35)
code(s, ML, y + 0.9, CW, 2.05, [
    ("# 1 \u00b7 start SQL Server 2022 in Docker", CODE_DIM),
    "docker run -d --name local-mssql \\",
    "  -e \"ACCEPT_EULA=Y\" \\",
    ("  -e \"MSSQL_SA_PASSWORD=<YOUR_PASSWORD>\" \\", CODE_KEY),
    "  -p 1433:1433 \\",
    "  mcr.microsoft.com/mssql/server:2022-latest",
], title="terminal")
code(s, ML, y + 3.28, CW, 0.92, [
    ("# 2 \u00b7 install ODBC Driver 17 for SQL Server on the host (needs admin / UAC)", CODE_DIM),
    ("learn.microsoft.com/en-us/sql/connect/odbc/download-odbc-driver-for-sql-server", CODE_KEY),
], title="host dependency")
callout(s, y + 4.42, "Give it ~30 seconds:",
        "SQL Server finishes initialising well after the container reports as running. Connecting "
        "too early looks exactly like a wrong password.", MS, h=0.56)
footer(s, "The walkthrough \u00b7 MS SQL Server")


# ============================== 9 · STEP 2 ==============================
s = slide()
y = head(s, 0.44, "Step 2 of 4 \u00b7 MS SQL Server", MS, "Seed a schema worth querying",
         "A deliberately different dataset from the other three databases", tsize=28)
txt(s, ML, y, CW, 0.72,
    [[("MS SQL Server gets a ", {}), ("books", {"font": C_FONT, "bold": True}),
      (" table rather than the ", {}), ("employees", {"font": C_FONT}),
      (" table used elsewhere. That is intentional: with all four connections live at once, "
       "distinct schemas make it immediately obvious which database actually answered a prompt.",
       {})]],
    size=13.5, spacing=1.35)
code(s, ML, y + 0.9, CW, 2.42, [
    ("-- sql/mssql/books_schema.sql", CODE_DIM),
    "CREATE TABLE books (",
    "  id     INT IDENTITY(1,1) PRIMARY KEY,",
    ("  title  VARCHAR(200) NOT NULL,", CODE_KEY),
    ("  author VARCHAR(100) NOT NULL,", CODE_KEY),
    "  genre  VARCHAR(50),",
    "  price  DECIMAL(6,2)",
    ");",
], title="schema")
txt(s, ML, y + 3.55, CW, 0.5,
    [[("Seeded with five rows \u2014 ", {"color": MUTED}),
      ("The Silent Orbit, Whispers of Kanto, Beneath the Ash, Letters to Aria, "
       "The Last Cartographer", {"color": TEXT, "italic": True}),
      (" \u2014 spanning five genres and a $11.99\u2013$18.00 price range, so aggregation and "
       "ranking queries both return something meaningful.", {"color": MUTED})]],
    size=12.5, spacing=1.3)
callout(s, y + 4.28, "Small on purpose:",
        "five rows exercise the whole path end to end. What matters is that the query is real, not "
        "that the dataset is large.", MS, h=0.56)
footer(s, "The walkthrough \u00b7 MS SQL Server")

# ============================== 10 · STEP 3 ==============================
s = slide()
y = head(s, 0.44, "Step 3 of 4 \u00b7 MS SQL Server", MS, "Install the MCP server, and pin the SDK",
         "The one step where the obvious command leaves you with a broken server", tsize=28)
txt(s, ML, y, CW, 0.95,
    [[("Install it with pip, not uvx", {"bold": True}),
      (" \u2014 the server needs to live in a real environment alongside the ODBC driver. Then "
       "immediately pin the SDK. The upstream ", {}),
      ("mcp", {"font": C_FONT}),
      (" package shipped a breaking 2.0.0 that renamed ", {}),
      ("FastMCP", {"font": C_FONT}),
      (" and moved the module entirely, but ", {}),
      ("sql-mcp-server", {"font": C_FONT}),
      (" declares no upper bound \u2014 so a clean install silently pulls the version that cannot "
       "run it.", {})]],
    size=13.5, spacing=1.35)
code(s, ML, y + 1.12, CW, 1.32, [
    "pip install sql-mcp-server",
    ("pip install \"mcp<2.0.0\" --force-reinstall", CODE_KEY),
    ("# second line is not optional \u2014 without it the server crashes on startup", CODE_DIM),
], title="terminal")
txt(s, ML, y + 2.62, CW, 0.3, "What you see if you skip the pin", size=10.5, bold=True, color=MUTED)
rect(s, ML, y + 2.9, CW, 0.62, fill=RGBColor(0xFE, 0xF2, 0xF2))
rect(s, ML, y + 2.9, 0.05, 0.62, fill=RED)
txt(s, ML + 0.26, y + 3.06, CW - 0.5, 0.35,
    [[("ModuleNotFoundError: No module named 'mcp.server.fastmcp'",
       {"font": C_FONT, "size": 12, "color": RED})]], spacing=1.2)
callout(s, y + 3.72, "Same bug, second time:",
        "PostgreSQL hit this identical breaking change. There it is pinned at invocation time inside "
        "the uvx command; here it has to be pinned at install time. Recognising it as one upstream "
        "problem rather than two saved a second round of debugging.", MS, h=0.78)
footer(s, "The walkthrough \u00b7 MS SQL Server")

# ============================== 11 · STEP 4 ==============================
s = slide()
y = head(s, 0.44, "Step 4 of 4 \u00b7 MS SQL Server", MS, "Register it with Claude Desktop",
         "The config block, and the two things that quietly go wrong here", tsize=28)
code(s, ML, y + 0.06, 7.15, 3.62, [
    "{",
    "  \"mcpServers\": {",
    ("    \"sql-assistant\": {", CODE_KEY),
    "      \"command\": \"sql-mcp-server\",",
    "      \"args\": [],",
    "      \"env\": {",
    "        \"DB_SERVER\": \"localhost\",",
    "        \"DB_NAME\": \"testdb\",",
    ("        \"DB_DRIVER\": \"ODBC Driver 17 for SQL Server\",", CODE_KEY),
    "        \"DB_USER\": \"sa\",",
    "        \"DB_PASSWORD\": \"<YOUR_PASSWORD>\"",
    "      } } } }",
], size=11, title="claude_desktop_config.json")

rx = ML + 7.45
rw = SW - MR - rx
txt(s, rx, y + 0.06, rw, 0.3, "Two traps in this step", size=12.5, bold=True, color=TEXT)
traps = [
    ("Find the config the right way", AMBER,
     "Use Settings \u2192 Developer \u2192 Edit Config. On Windows, packaged installs of Claude Desktop "
     "do not use %APPDATA%\\Claude \u2014 the real file sits under AppData\\Local\\Packages\\<id>\\"
     "LocalCache\\Roaming\\Claude\\. Editing the wrong copy looks like the config being ignored."),
    ("Quit properly, not just close", RED,
     "Closing the window leaves it running in the tray, and MCP servers are only launched on a "
     "cold start. Quit from the system tray, reopen, then confirm sql-assistant reads running "
     "under Settings \u2192 Developer."),
]
ty = y + 0.46
for t, c, b in traps:
    rect(s, rx, ty, rw, 1.58, fill=PANEL)
    rect(s, rx, ty, 0.05, 1.58, fill=c)
    txt(s, rx + 0.24, ty + 0.16, rw - 0.46, 0.26, t, size=12.5, bold=True, color=TEXT)
    txt(s, rx + 0.24, ty + 0.5, rw - 0.46, 1.0, b, size=11, color=MUTED, spacing=1.26)
    ty += 1.72
callout(s, y + 3.88, "Add, don't replace:",
        "all four servers coexist under the same mcpServers key \u2014 keep the existing entries and "
        "append this one.", MS, h=0.56)
footer(s, "The walkthrough \u00b7 MS SQL Server")


# ============================== 12\u201315 · VALIDATION ==============================
def two_panel(s, pa, pb, y0, avail_h, gap=0.2):
    """Two crops at one shared scale, side by side, so text size matches across both."""
    with Image.open(pa) as ia:
        wa, ha = ia.size
    with Image.open(pb) as ib:
        wb, hb = ib.size
    colw = (CW - gap) / 2
    sc = min(colw / wa, colw / wb, avail_h / ha, avail_h / hb)
    Wa, Ha, Wb, Hb = wa * sc, ha * sc, wb * sc, hb * sc
    x = ML + (CW - (Wa + gap + Wb)) / 2
    y = y0 + (avail_h - max(Ha, Hb)) / 2          # centre when width-bound
    for p, px, pw, ph in ((pa, x, Wa, Ha), (pb, x + Wa + gap, Wb, Hb)):
        pic = s.shapes.add_picture(p, Inches(px), Inches(y), Inches(pw), Inches(ph))
        try:
            pic.line.color.rgb = RULE
            pic.line.width = Pt(0.75)
        except Exception:
            pass
    return sc


def showcase(stem, kicker, title, lead, prompt):
    s = slide()
    chip(s, 0.34, kicker, MS)
    txt(s, ML, 0.62, 5.05, 0.4, title, size=21, bold=True, color=TEXT, spacing=0.98)
    txt(s, ML, 1.0, 5.05, 0.42, lead, size=11.5, color=MUTED, spacing=1.26)
    px = ML + 5.35
    pw = SW - MR - px
    txt(s, px, 0.36, pw, 0.24, "ASKED IN CLAUDE DESKTOP", size=9, bold=True, color=FAINT)
    rect(s, px, 0.62, pw, 0.72, fill=PANEL)
    rect(s, px, 0.62, 0.05, 0.72, fill=BLUE)
    txt(s, px + 0.24, 0.78, pw - 0.44, 0.5,
        [[("\u201c" + prompt + "\u201d", {"italic": True, "size": 12.5, "color": TEXT})]],
        spacing=1.26)
    two_panel(s, f"docs/deck-assets/{stem}_a.png", f"docs/deck-assets/{stem}_b.png", 1.5, 5.42)
    footer(s, "MS SQL Server")


showcase(
    "mssql_query_schema",
    "MS SQL Server \u00b7 schema discovery",
    "Reading the schema",
    "No schema is supplied in the prompt \u2014 the column names, types and identity key all come back "
    "from a live metadata query.",
    "What columns does the books table have?")

showcase(
    "mssql_query_aggregation",
    "MS SQL Server \u00b7 aggregation",
    "Aggregating across the table",
    "Picks the grouping and the measure itself, composes the GROUP BY, and returns per-genre "
    "averages computed from the actual rows.",
    "What's the average price by genre in the books table?")

showcase(
    "mssql_query_topquery",
    "MS SQL Server \u00b7 query generation",
    "Writing dialect-specific T-SQL",
    "Uses TOP rather than the LIMIT syntax the other three engines expect \u2014 the same prompt "
    "produces different, correct SQL per database.",
    "Write and run a query to find the most expensive book.")

showcase(
    "mssql_query_codegen",
    "MS SQL Server \u00b7 code generation",
    "Writing code from live data",
    "Queries the table first, then writes a script shaped around the schema it just read \u2014 real "
    "column names and types, no placeholders to fix afterwards.",
    "Write a Python script that pulls this data and charts it.")


# ============================== 16 · WHAT BROKE ==============================
s = slide()
y = head(s, 0.46, "Lessons", RED, "What broke, and how it was found",
         "None of this was in the original guide \u2014 each layer had to be diagnosed independently")
gw = (CW - 0.3) / 2
gh = 1.86
faults = [
    ("The recommended image spoke the wrong protocol", RED,
     "openmcpserver/mcp-postgres:latest runs an HTTP server under Uvicorn. Claude Desktop's local "
     "MCP servers communicate over stdio, so this image cannot work here at all \u2014 confirmed by "
     "running it manually and watching it bind an HTTP port instead of starting a stdio process."),
    ("The officially recommended package was deprecated", AMBER,
     "@modelcontextprotocol/server-postgres was retired mid-2026. Its suggested replacement, "
     "@crystaldba/postgres-mcp, does not exist on npm at all \u2014 it is a PyPI package meant to be "
     "run with uvx, not npx, which is why every npx attempt failed."),
    ("An upstream SDK shipped a breaking change mid-build", MS,
     "mcp 2.0.0 renamed FastMCP and moved the module. Neither postgres-mcp nor sql-mcp-server "
     "declares an upper bound, so both resolved the version that breaks them. Pinning mcp<2.0.0 "
     "fixed both \u2014 at invocation for uvx, at install for pip."),
    ("A native Windows service was hijacking the port", PG,
     "Postgres auth kept failing with a provably correct password. netstat -ano | findstr :5432 "
     "showed two listening PIDs; tasklist identified one as postgres.exe running as a Windows "
     "service, intercepting connections before Docker ever saw them. The container moved to 5433."),
]
for i, (t, c, b) in enumerate(faults):
    gx = ML + (i % 2) * (gw + 0.3)
    gy = y + 0.1 + (i // 2) * (gh + 0.18)
    rect(s, gx, gy, gw, gh, fill=PANEL)
    rect(s, gx, gy, gw, 0.05, fill=c)
    txt(s, gx + 0.26, gy + 0.24, gw - 0.5, 0.48, t, size=12.5, bold=True, color=TEXT, spacing=1.1)
    txt(s, gx + 0.26, gy + 0.78, gw - 0.5, 1.0, b, size=10.5, color=MUTED, spacing=1.26)
callout(s, y + 0.1 + 2 * gh + 0.18 + 0.26, "The pattern:",
        "every failure surfaced far from its cause \u2014 a driver problem read as a bad password, a "
        "port conflict read as bad credentials, a dependency break read as a broken package.",
        RED, h=0.56)
footer(s, "Lessons")

# ============================== 17 · LIMITS ==============================
s = slide()
y = head(s, 0.46, "Lessons", RED, "Being straight about the limits",
         "What this setup does not do, stated plainly")
lims = [
    ("Read-only is a wrapper flag, not a database rule", AMBER,
     "POSTGRES_READ_ONLY is enforced by the MCP wrapper \u2014 the underlying user is still a "
     "superuser. Anything bypassing the wrapper could write. A rigorous setup would create a "
     "dedicated read-only role and GRANT SELECT instead of relying on the flag."),
    ("Nothing persists across a reboot", BLUE,
     "Docker containers do not auto-start, and Claude Desktop only launches MCP servers on a cold "
     "start. Bringing everything back up is a documented five-minute sequence, not an "
     "always-on service."),
    ("The datasets are deliberately tiny", GREEN,
     "A handful of rows per database. The aim was a working protocol path end to end, not a "
     "benchmark of query performance or data volume."),
    ("No real credentials are committed", MS,
     "Config files ship with placeholders, secrets live in a gitignored .env alongside a checked-in "
     ".env.example, and the real Claude Desktop config is gitignored too."),
]
ly = y + 0.12
for t, c, b in lims:
    rect(s, ML, ly, 0.05, 0.92, fill=c)
    txt(s, ML + 0.26, ly, CW - 0.4, 0.26, t, size=13, bold=True, color=TEXT)
    txt(s, ML + 0.26, ly + 0.32, CW - 0.4, 0.62, b, size=11, color=MUTED, spacing=1.26)
    ly += 1.04
callout(s, ly + 0.06, "Why list these at all:",
        "knowing the limits of a control is more useful than claiming it is airtight \u2014 and every "
        "one of these is written up in the repo rather than left implied.", RED, h=0.56)
footer(s, "Lessons")

# ============================== 18 · EXPLORE ==============================
s = slide(dark=True)
rect(s, 0, 0, 0.16, SH, fill=BLUE)
txt(s, ML, 0.72, 3.0, 0.35, "EXPLORE IT", size=12, bold=True, color=BLUE)
txt(s, ML, 1.12, 8.6, 0.9, "Everything is in the repo", size=34, bold=True, color=WHITE)
txt(s, ML, 2.08, 7.2, 0.9,
    "Each database has its own guide with every command, config block and gotcha for that engine \u2014 "
    "this deck walked through one of the four.",
    size=14, color=RGBColor(0xB6, 0xBD, 0xC8), spacing=1.35)
rect(s, ML, 3.05, 7.05, 0.72, fill=INK_SOFT)
txt(s, ML + 0.28, 3.24, 6.6, 0.4,
    [[("github.com/4reeb-5yed/claude-db-connect", {"font": C_FONT, "size": 15, "bold": True,
                                                   "color": WHITE})]], spacing=1.2)
guides = [("docs/setup/postgres.md", PG), ("docs/setup/mysql.md", MY),
          ("docs/setup/mssql.md", MS), ("docs/setup/sqlite.md", SQ)]
for i, (g, c) in enumerate(guides):
    gy = 4.08 + i * 0.56
    rect(s, ML, gy + 0.09, 0.16, 0.16, fill=c)
    txt(s, ML + 0.34, gy, 6.6, 0.35,
        [[(g, {"font": C_FONT, "size": 12.5, "color": RGBColor(0xD5, 0xDA, 0xE2)})]], spacing=1.2)
rect(s, 8.55, 3.05, 4.05, 3.35, fill=INK_SOFT)
rect(s, 8.55, 3.05, 0.05, 3.35, fill=BLUE)
txt(s, 8.85, 3.3, 3.5, 0.3, "Also in the repo", size=12, bold=True, color=BLUE)
extras = [("README.md", "overview, architecture and a screenshot per database"),
          ("PLAN.md", "the full build log and design decisions"),
          ("scripts/", "a helper to bring the Postgres container up"),
          ("sql/ \u00b7 config/", "schemas and Claude Desktop config examples")]
ey = 3.75
for t, d in extras:
    txt(s, 8.85, ey, 3.5, 0.28,
        [[(t, {"font": C_FONT, "size": 11.5, "bold": True, "color": WHITE})]], spacing=1.15)
    txt(s, 8.85, ey + 0.26, 3.5, 0.4, d, size=10.5, color=RGBColor(0x8E, 0x96, 0xA3), spacing=1.2)
    ey += 0.66
txt(s, ML, 6.62, 11.5, 0.4,
    [[("Four databases, four different failure modes, one working protocol path \u2014 all documented.",
       {"italic": True, "size": 12.5, "color": RGBColor(0x8E, 0x96, 0xA3)})]])
footer(s, "Claude DB Connect", dark=True)

prs.save(OUT)
print(f"OK  {OUT}  slides={len(prs.slides.__iter__.__self__._sldIdLst)}")
